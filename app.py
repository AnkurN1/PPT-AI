import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import os

# Load data
data = pd.read_excel("all companys database.xlsx")
IMAGE_BASE = "images"
LOGO_BASE = "logo"

# Utility functions
def clean_string(s):
    return str(s).replace("\xa0", " ").replace("\n", "").strip()

def get_image_list(company, product, ptype):
    folder = os.path.join(IMAGE_BASE, str(company).strip(), str(product).strip(), str(ptype).strip())
    images = []
    if os.path.exists(folder):
        for file in os.listdir(folder):
            if file.lower().endswith(('.jpg', '.jpeg', '.png')):
                images.append(os.path.join(folder, file))
    return images

def get_scaled_dimensions(img, max_width, max_height):
    img_width_px, img_height_px = img.size
    aspect_ratio = img_width_px / img_height_px
    box_aspect = max_width / max_height

    if aspect_ratio > box_aspect:
        width = max_width
        height = width / aspect_ratio
    else:
        height = max_height
        width = height * aspect_ratio

    return width, height

def create_beautiful_ppt(slides, include_intro_outro=True):
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    first_slide_path = "img/first.png"
    last_slide_path = "img/last.png"

    # First slide
    if include_intro_outro and os.path.exists(first_slide_path):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(first_slide_path, Inches(0), Inches(0),
                                 width=prs.slide_width, height=prs.slide_height)

    # Product slides
    for slide_data in slides:
        slide = prs.slides.add_slide(blank)
        company = slide_data.get('company', '')

        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12), Inches(0.7))
        frame = title_shape.text_frame
        frame.text = slide_data['title']
        p = frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0, 0, 0)

        # Images
        y_img_top = 1.2
        y_img_bottom = 6.9
        max_img_height = y_img_bottom - y_img_top
        slide_width_in = prs.slide_width.inches
        imgs = slide_data['images']
        img_count = len(imgs)

        if img_count == 1:
            img_path = imgs[0]
            with Image.open(img_path) as img:
                img_width, img_height = get_scaled_dimensions(img, max_width=10.5, max_height=max_img_height)
                x = (slide_width_in - img_width) / 2
                y = y_img_top + (max_img_height - img_height) / 2
                slide.shapes.add_picture(img_path, Inches(x), Inches(y),
                                         width=Inches(img_width), height=Inches(img_height))
        else:
            padding = 0.2
            columns = min(img_count, 3)
            rows = (img_count + columns - 1) // columns
            available_width = slide_width_in - (padding * (columns + 1))
            available_height = max_img_height - ((rows - 1) * padding)
            cell_width = available_width / columns
            cell_height = available_height / rows

            for i, img_path in enumerate(imgs):
                row = i // columns
                col = i % columns
                with Image.open(img_path) as img:
                    img_width, img_height = get_scaled_dimensions(img, max_width=cell_width, max_height=cell_height)
                    x = padding + col * (cell_width + padding) + (cell_width - img_width) / 2
                    y = y_img_top + row * (cell_height + padding) + (cell_height - img_height) / 2
                    slide.shapes.add_picture(img_path, Inches(x), Inches(y),
                                             width=Inches(img_width), height=Inches(img_height))

        # Logo (top-right)
        logo_path = os.path.join(LOGO_BASE, company, "logo.png")
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, prs.slide_width - Inches(1.2), Inches(0.1), width=Inches(1.1))

        # Copyright (bottom-right)
        cp_box = slide.shapes.add_textbox(prs.slide_width - Inches(3.6), prs.slide_height - Inches(0.3), Inches(3.6), Inches(0.4))
        cp_frame = cp_box.text_frame
        cp_frame.text = "Copyright © 2025 Altossa Projects LLp. All Rights Reserved."
        cp_para = cp_frame.paragraphs[0]
        cp_para.font.size = Pt(10)
        cp_para.font.color.rgb = RGBColor(128, 128, 128)

        # Link (bottom-left)
        if slide_data.get('link'):
            link_box = slide.shapes.add_textbox(Inches(0.1), prs.slide_height - Inches(0.3), Inches(7), Inches(0.4))
            link_frame = link_box.text_frame
            link_frame.text = str(slide_data.get('link') or "")
            p = link_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0, 102, 204)

    # Last slide
    if include_intro_outro and os.path.exists(last_slide_path):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(last_slide_path, Inches(0), Inches(0),
                                 width=prs.slide_width, height=prs.slide_height)

    output_path = "styled_ppt.pptx"
    prs.save(output_path)
    return output_path

# Streamlit App
st.title("Product Type Search & Select")

search_query = st.text_input("Search by Type", "")

search_selected_items = []
manual_selected_items = []

# Search View
if search_query:
    filtered_data = data[data['Type'].str.contains(search_query, case=False, na=False)]
    for idx, row in filtered_data.iterrows():
        product = row['Product']
        ptype = row['Type']
        link = row.get("Link", "")
        company = row['Company']
        img_paths = get_image_list(company, product, ptype)

        st.markdown(f"### {product} - {ptype}")
        img_cols = st.columns(min(4, len(img_paths)))
        selected_imgs = []

        for i, path in enumerate(img_paths):
            with img_cols[i % len(img_cols)]:
                st.image(path, width=160)
                key = f"search_{company}_{product}_{ptype}_{i}".replace(" ", "_")
                if st.checkbox("Include", key=key):
                    selected_imgs.append(path)

        if selected_imgs:
            search_selected_items.append({
                "title": f"{product}",
                "company": company,
                "link": link,
                "images": selected_imgs
            })

    if st.button("Generate for Search Results"):
        if not search_selected_items:
            st.warning("Please select at least one image.")
        else:
            ppt_path = create_beautiful_ppt(search_selected_items, include_intro_outro=False)
            with open(ppt_path, "rb") as f:
                st.download_button("Download PPT", f, file_name="search_presentation.pptx")

# Manual View
else:
    st.markdown("## Select Manually")

    company = st.selectbox("Select Company", sorted(data['Company'].dropna().unique()), key="select_company_dropdown")
    products = sorted(data[data['Company'] == company]['Product'].dropna().unique())
    product = st.selectbox("Select Product", products, key="select_product_dropdown")
    filtered_rows = data[(data['Company'] == company) & (data['Product'] == product)]

    for idx, row in filtered_rows.iterrows():
        ptype = row['Type']
        link = row.get("Link", "")
        img_paths = get_image_list(company, product, ptype)

        st.markdown(f"### {product} - {ptype}")
        img_cols = st.columns(min(4, len(img_paths)))
        selected_imgs = []

        for i, path in enumerate(img_paths):
            with img_cols[i % len(img_cols)]:
                st.image(path, width=160)
                key = f"manual_{company}_{product}_{ptype}_{i}".replace(" ", "_")
                if st.checkbox("Include", key=key):
                    selected_imgs.append(path)

        if selected_imgs:
            manual_selected_items.append({
                "title": f"{product}",
                "company": company,
                "link": link,
                "images": selected_imgs
            })

    if st.button("Generate PPT"):
        if not manual_selected_items:
            st.warning("Please select at least one image.")
        else:
            ppt_path = create_beautiful_ppt(manual_selected_items, include_intro_outro=True)
            with open(ppt_path, "rb") as f:
                st.download_button("Download PPT", f, file_name=f"{company}_{product}_presentation.pptx")

